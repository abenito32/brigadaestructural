#!/usr/bin/env python3
# Brigada · Evaluación estructural en campo
# Copyright (C) 2026 Rollout Comercio e Servicios Limitada / Andrés Benito Revollo Vélez
# Software libre bajo GNU AGPL v3.
"""
Genera la presentación institucional (.pptx) para alcaldías, universidades y
asociaciones profesionales.

  pip install python-pptx
  python3 docs/build_presentacion.py

El .pptx queda editable: quien lo reciba o quien lo presente puede ajustar textos
sin volver a correr esto.
"""
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

RAIZ = pathlib.Path(__file__).resolve().parent
SALIDA = RAIZ / "brigada-presentacion.pptx"
ASSETS = RAIZ / "assets"

# Paleta de la aplicación, para que la presentación y el producto se parezcan.
AZUL = RGBColor(0x03, 0x69, 0xA1)
AZUL_OSC = RGBColor(0x07, 0x59, 0x85)
TINTA = RGBColor(0x0F, 0x17, 0x2A)
TINTA2 = RGBColor(0x47, 0x55, 0x69)
TENUE = RGBColor(0x5E, 0x6E, 0x82)
PAPEL = RGBColor(0xF1, 0xF5, 0xF9)
LINEA = RGBColor(0xD9, 0xE0, 0xE8)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
VERDE = RGBColor(0x15, 0x80, 0x3D)
AMBAR = RGBColor(0xFA, 0xCC, 0x15)
AMBAR_T = RGBColor(0x42, 0x20, 0x06)
NARANJA = RGBColor(0xC2, 0x41, 0x0C)
ROJO = RGBColor(0x7F, 0x1D, 0x1D)

FUENTE = "Arial"          # disponible en Windows, macOS y LibreOffice
ANCHO, ALTO = Inches(13.333), Inches(7.5)

# El correo NO va en el código: este archivo está en un repositorio público y una
# dirección ahí se rastrea en días. Se pasa por entorno al generar:
#   BRIGADA_CONTACTO=correo@dominio python3 docs/build_presentacion.py
import os

CONTACTO = {
    "nombre": "Andrés Benito Revollo Vélez",
    "empresa": "Rollout Comercio e Servicios Limitada",
    "web": "www.artificialintelligence.tec.br",
    "correo": os.getenv("BRIGADA_CONTACTO", "—"),
}


# ─────────────────────────────────────────────────────────── utilidades
def texto(dia, x, y, w, h, contenido, *, tam=18, color=TINTA, negrita=False,
          align=PP_ALIGN.LEFT, interlineado=1.25, anchor=MSO_ANCHOR.TOP):
    caja = dia.shapes.add_textbox(x, y, w, h)
    tf = caja.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lineas = contenido if isinstance(contenido, list) else [contenido]
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = interlineado
        if isinstance(linea, tuple):          # (texto, {overrides})
            linea, extra = linea
        else:
            extra = {}
        r = p.add_run()
        r.text = linea
        f = r.font
        f.name, f.size = FUENTE, Pt(extra.get("tam", tam))
        f.bold = extra.get("negrita", negrita)
        f.color.rgb = extra.get("color", color)
        if extra.get("espacio_antes"):
            p.space_before = Pt(extra["espacio_antes"])
    return caja


def caja(dia, x, y, w, h, *, relleno=BLANCO, borde=LINEA, radio=True, grosor=1):
    forma = dia.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radio else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radio and forma.adjustments:      # el rectangulo recto no tiene ajuste
        forma.adjustments[0] = 0.045
    forma.fill.solid()
    forma.fill.fore_color.rgb = relleno
    if borde is None:
        forma.line.fill.background()
    else:
        forma.line.color.rgb = borde
        forma.line.width = Pt(grosor)
    forma.shadow.inherit = False
    forma.text_frame.text = ""
    return forma


def barra_semaforo(dia, x, y, ancho_total, alto=Inches(0.22)):
    """El motivo de la portada: las tres franjas del semáforo."""
    props = [(VERDE, .22), (AMBAR, .26), (NARANJA, .28), (ROJO, .24)]
    cursor = x
    for color, frac in props:
        w = Emu(int(ancho_total * frac) - Inches(0.08))
        f = caja(dia, cursor, y, w, alto, relleno=color, borde=None)
        cursor = Emu(cursor + w + Inches(0.08))
    return


def nueva(prs, *, fondo=BLANCO):
    dia = prs.slides.add_slide(prs.slide_layouts[6])   # en blanco
    f = dia.background.fill
    f.solid()
    f.fore_color.rgb = fondo
    return dia


def cabecera(dia, rotulo, titulo, bajada=None):
    """Devuelve la Y donde puede empezar el contenido, para no adivinarla."""
    texto(dia, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.3), rotulo,
          tam=11, color=AZUL, negrita=True)
    # 28pt y no 32: a 32 varios titulos pasaban a dos lineas y pisaban la barra.
    texto(dia, Inches(0.9), Inches(0.92), Inches(11.5), Inches(0.55), titulo,
          tam=28, color=TINTA, negrita=True)
    caja(dia, Inches(0.9), Inches(1.62), Inches(1.2), Pt(3.5), relleno=AZUL,
         borde=None, radio=False)
    if bajada:
        texto(dia, Inches(0.9), Inches(1.9), Inches(11.5), Inches(0.4), bajada,
              tam=14, color=TINTA2)
        return Inches(2.55)
    return Inches(2.2)


def pie_contacto(dia, color=TENUE):
    texto(dia, Inches(0.9), Inches(6.92), Inches(11.5), Inches(0.3),
          f"{CONTACTO['empresa']}  ·  {CONTACTO['web']}  ·  {CONTACTO['correo']}",
          tam=10, color=color)


def tarjetas(dia, items, *, cols=3, y=Inches(2.55), alto=Inches(1.95),
             x0=Inches(0.9), ancho_total=Inches(11.5), gap=Inches(0.28)):
    w = Emu(int((ancho_total - gap * (cols - 1)) / cols))
    for i, (titulo, cuerpo) in enumerate(items):
        fila, col = divmod(i, cols)
        x = Emu(x0 + col * (w + gap))
        yy = Emu(y + fila * (alto + gap))
        caja(dia, x, yy, w, alto, relleno=PAPEL, borde=LINEA)
        # El titulo tiene sitio para dos lineas: los renderizadores no miden igual.
        texto(dia, Emu(x + Inches(0.3)), Emu(yy + Inches(0.24)),
              Emu(w - Inches(0.6)), Inches(0.6), titulo, tam=14, negrita=True,
              color=AZUL_OSC, interlineado=1.15)
        texto(dia, Emu(x + Inches(0.3)), Emu(yy + Inches(0.92)),
              Emu(w - Inches(0.6)), Emu(alto - Inches(1.15)), cuerpo, tam=11.5,
              color=TINTA2, interlineado=1.25)


# ─────────────────────────────────────────────────────────── diapositivas
def portada(prs):
    dia = nueva(prs, fondo=PAPEL)
    caja(dia, Emu(0), Emu(0), Inches(0.34), ALTO, relleno=AZUL, borde=None, radio=False)
    icono = ASSETS / "icono-512.png"
    if icono.exists():
        dia.shapes.add_picture(str(icono), Inches(10.4), Inches(0.85), Inches(2.0), Inches(2.0))
    texto(dia, Inches(1.1), Inches(1.15), Inches(9.0), Inches(0.35),
          "EVALUACIÓN ESTRUCTURAL POST-SISMO", tam=13, color=AZUL, negrita=True)
    texto(dia, Inches(1.1), Inches(1.7), Inches(8.9), Inches(0.9),
          "Brigada", tam=54, color=TINTA, negrita=True, interlineado=1.0)
    texto(dia, Inches(1.1), Inches(2.72), Inches(8.9), Inches(1.1),
          ["Levantamiento de daño en campo,", "sin depender de la señal"],
          tam=30, color=TINTA2, interlineado=1.2)
    barra_semaforo(dia, Inches(1.1), Inches(4.3), Inches(5.2))
    texto(dia, Inches(1.1), Inches(4.8), Inches(9.5), Inches(0.9),
          ["Aplicación para que ingenieros con matrícula evalúen edificaciones tras un sismo",
           "y consoliden los resultados de forma trazable y verificable."],
          tam=16, color=TINTA2, interlineado=1.35)
    caja(dia, Inches(1.1), Inches(5.85), Inches(11.2), Pt(1), relleno=LINEA, borde=None, radio=False)
    texto(dia, Inches(1.1), Inches(6.1), Inches(6.5), Inches(1.0),
          [(CONTACTO["nombre"], {"negrita": True, "color": TINTA, "tam": 14}),
           (CONTACTO["empresa"], {"tam": 12, "color": TINTA2})], tam=12, color=TINTA2)
    texto(dia, Inches(7.6), Inches(6.1), Inches(4.7), Inches(1.0),
          [(CONTACTO["web"], {"tam": 12, "color": AZUL_OSC, "negrita": True}),
           (CONTACTO["correo"], {"tam": 12, "color": AZUL_OSC})],
          align=PP_ALIGN.RIGHT)
    return dia


def problema(prs):
    dia = nueva(prs)
    cabecera(dia, "EL PROBLEMA", "Después del sismo, la información es el cuello de botella",
             "Las primeras 72 horas deciden a dónde va la ayuda. Hoy esa información se levanta en papel.")
    tarjetas(dia, [
        ("Sin señal, no hay sistema",
         "Las herramientas en línea dejan de servir justo donde se necesitan. La red celular "
         "es lo primero que cae y lo último que vuelve."),
        ("El papel se pierde",
         "Formatos que se mojan, se traspapelan o llegan días después. Digitalizarlos consume "
         "el tiempo del personal que hace falta en terreno."),
        ("Criterios que no coinciden",
         "Cada brigada evalúa con su propio criterio. Consolidar veinte planillas distintas "
         "produce un número en el que nadie confía."),
        ("Nadie sabe quién firmó qué",
         "Sin trazabilidad no hay forma de auditar una decisión que ordenó evacuar un edificio "
         "o que permitió volver a habitarlo."),
        ("Duplicados y vacíos",
         "Manzanas evaluadas tres veces mientras otras quedan sin cubrir, porque no hay una "
         "vista común del avance."),
        ("El dato personal se filtra",
         "Direcciones de predios circulando por grupos de mensajería. Además del riesgo para "
         "las personas, es un incumplimiento de la Ley 1581 de 2012."),
    ])
    pie_contacto(dia)
    return dia


def que_es(prs):
    dia = nueva(prs)
    cabecera(dia, "LA PROPUESTA", "Una herramienta de triaje que funciona donde no hay red")
    izq = Inches(0.9)
    texto(dia, izq, Inches(2.3), Inches(6.4), Inches(3.1), [
        ("Se instala en el teléfono y funciona completa sin conexión.",
         {"tam": 16, "negrita": True, "color": TINTA}),
        ("Todo se guarda en el equipo con almacenamiento duradero —el navegador no lo "
         "desaloja para hacer espacio— y se sincroniza cuando vuelve la señal. La cola "
         "sobrevive al cierre de la aplicación y al reinicio del teléfono.",
         {"tam": 13, "color": TINTA2, "espacio_antes": 9}),
        ("Aplica el semáforo ATC-20 adaptado a la NSR-10: se calcula solo a partir del "
         "daño observado, y el ingeniero puede modificarlo dejando constancia escrita.",
         {"tam": 13, "color": TINTA2, "espacio_antes": 9}),
        ("No es auto-reporte ciudadano: quien evalúa y firma es un profesional con "
         "matrícula vigente, en sitio.",
         {"tam": 13, "color": TINTA2, "espacio_antes": 9}),
    ], interlineado=1.3)
    caja(dia, izq, Inches(5.5), Inches(6.4), Inches(1.05), relleno=PAPEL, borde=LINEA)
    texto(dia, Emu(izq + Inches(0.3)), Inches(5.72), Inches(5.8), Inches(0.7),
          ["Es triaje preliminar. La habilitación definitiva de una edificación es",
           "competencia de UNGRD, Defensa Civil, bomberos y las alcaldías."],
          tam=12, color=TINTA2, interlineado=1.3)
    captura = ASSETS / "app-escala.png"
    if captura.exists():
        caja(dia, Inches(7.75), Inches(2.3), Inches(4.65), Inches(4.25),
             relleno=PAPEL, borde=LINEA)
        dia.shapes.add_picture(str(captura), Inches(7.9), Inches(2.45), width=Inches(4.35))
    pie_contacto(dia)
    return dia


def como_funciona(prs):
    dia = nueva(prs)
    cabecera(dia, "CÓMO FUNCIONA", "Tres pasos, y ninguno depende de que haya red")
    pasos = [
        ("1", "Se prepara una vez", "Con señal, el inspector abre la dirección, instala la "
         "aplicación en su teléfono y registra su nombre, matrícula y brigada. Cinco minutos, "
         "una sola vez por equipo."),
        # Cabe en cinco líneas dentro de la tarjeta. Al agregar la ruta hubo que soltar
        # «en cuatro categorías»: con seis líneas el texto se salía del recuadro.
        ("2", "Se evalúa sin conexión", "En terreno: ubicación, tipo de edificación, nivel de "
         "daño, condiciones de cierre y fotos. Con ruta asignada, parada por parada. "
         "El semáforo se calcula solo."),
        ("3", "Se sincroniza al volver", "Un botón sube todo lo pendiente y baja la ruta del "
         "día. El envío es idempotente: reintentar nunca duplica, y el servidor devuelve un "
         "acuse que queda guardado en el teléfono."),
    ]
    x0, w, gap = Inches(0.9), Inches(3.65), Inches(0.28)
    for i, (n, titulo, cuerpo) in enumerate(pasos):
        x = Emu(x0 + i * (w + gap))
        caja(dia, x, Inches(2.5), w, Inches(3.1), relleno=BLANCO, borde=LINEA)
        circ = dia.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + Inches(0.32)), Inches(2.85),
                                    Inches(0.62), Inches(0.62))
        circ.fill.solid(); circ.fill.fore_color.rgb = AZUL
        circ.line.fill.background(); circ.shadow.inherit = False
        tf = circ.text_frame; tf.text = n
        r = tf.paragraphs[0].runs[0]
        r.font.name, r.font.size, r.font.bold = FUENTE, Pt(20), True
        r.font.color.rgb = BLANCO
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        texto(dia, Emu(x + Inches(0.32)), Inches(3.68), Emu(w - Inches(0.64)), Inches(0.4),
              titulo, tam=17, negrita=True, color=TINTA)
        texto(dia, Emu(x + Inches(0.32)), Inches(4.12), Emu(w - Inches(0.64)), Inches(1.2),
              cuerpo, tam=12.5, color=TINTA2, interlineado=1.3)
    caja(dia, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.85), relleno=PAPEL, borde=LINEA)
    texto(dia, Inches(1.2), Inches(6.0), Inches(10.9), Inches(0.4),
          "El servidor responde 503 y nunca 200 si no pudo grabar: así la aplicación conserva "
          "el pendiente en vez de darlo por enviado. Ninguna evaluación se pierde en silencio.",
          tam=13, color=TINTA2)
    pie_contacto(dia)
    return dia


def semaforo(prs):
    dia = nueva(prs)
    cabecera(dia, "CRITERIO TÉCNICO", "Los cuatro niveles del formulario V2F del IDIGER",
             "La regla es determinista: los mismos datos producen siempre la misma clasificación.")
    filas = [
        (ROJO, BLANCO, "4 · COLAPSO", "Peligro de colapso",
         "Evacuar y acordonar el entorno. Colapso, inclinación visible, o núcleo de "
         "columna triturado."),
        (NARANJA, BLANCO, "3 · NO HAB.", "No habitable",
         "Evacuar. Daño severo en portantes o entrepisos, grietas pasantes, riesgo "
         "externo."),
        (AMBAR, AMBAR_T, "2 · RESTRING.", "Uso restringido",
         "Daño moderado en portantes, entrepisos o terreno; o severo solo en muros "
         "divisorios y fachada."),
        (VERDE, BLANCO, "1 · HABITABLE", "Habitable",
         "Daño leve o nulo, sin ninguna condición de cierre presente."),
    ]
    y = Inches(2.35)
    for fondo, tinta_chip, nombre, sentido, regla in filas:
        caja(dia, Inches(0.9), y, Inches(7.6), Inches(0.92), relleno=BLANCO, borde=LINEA)
        caja(dia, Inches(1.15), Emu(y + Inches(0.19)), Inches(1.6), Inches(0.52),
             relleno=fondo, borde=None)
        texto(dia, Inches(1.15), Emu(y + Inches(0.31)), Inches(1.6), Inches(0.3), nombre,
              tam=11.5, negrita=True, color=tinta_chip, align=PP_ALIGN.CENTER)
        texto(dia, Inches(2.95), Emu(y + Inches(0.14)), Inches(5.3), Inches(0.3), sentido,
              tam=13, negrita=True, color=TINTA)
        texto(dia, Inches(2.95), Emu(y + Inches(0.42)), Inches(5.3), Inches(0.46), regla,
              tam=10.5, color=TINTA2, interlineado=1.12)
        y = Emu(y + Inches(1.02))
    captura = ASSETS / "app-placa.png"
    if captura.exists():
        caja(dia, Inches(8.85), Inches(2.35), Inches(3.55), Inches(3.9),
             relleno=PAPEL, borde=LINEA)
        dia.shapes.add_picture(str(captura), Inches(9.0), Inches(2.5), width=Inches(3.25))
    texto(dia, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.36),
          "La global es la peor de cinco parciales —estado general, geotécnicos, no "
          "estructurales, estructurales y entorno—, como manda el V2F.",
          tam=11, color=TENUE, interlineado=1.2)
    pie_contacto(dia)
    return dia


def panel(prs):
    """Lo que ve quien coordina. Es la mitad que no se ve en una demo del teléfono."""
    dia = nueva(prs)
    cabecera(dia, "COORDINACIÓN", "Quien dirige la operación ve lo que llega, en vivo",
             "Panel web en dirección separada. Los inspectores no tienen cuenta: "
             "su matrícula es una firma, no un acceso.")
    tarjetas(dia, [
        ("Mapa del consolidado",
         "Un círculo por sector sobre callejero, satélite o híbrida. El área es cuántas "
         "evaluaciones tiene; el color, qué proporción quedó sin poder habitarse."),
        ("El reporte completo, con fotos",
         "Cada evaluación se abre con su registro fotográfico y se exporta en un archivo "
         "que se abre sin conexión, para adjuntarlo a un oficio."),
        ("Doble revisión de los desalojos",
         "No habitable y peligro de colapso entran en cola; el segundo va primero y con "
         "plazo más corto. Quien firmó no puede revisar lo suyo."),
        ("Evolución de la operación",
         "Ritmo por día, qué sectores llevan horas sin actividad y qué queda por resolver. "
         "Sin conteo por persona: la prisa es el riesgo."),
        ("Cada brigada ve lo suyo",
         "El coordinador de una brigada entra con su propia clave y ve solo sus reportes. "
         "El alcance se decide en el servidor, no escondiendo un menú."),
        ("Exportación al formulario V2F",
         "Una columna por casilla del formulario del IDIGER, con sus códigos, en CSV, "
         "JSON o por API. Sin transcribir nada a mano."),
    ])
    pie_contacto(dia)
    return dia


def despacho(prs):
    """Lo que convierte «se hicieron 40 evaluaciones» en «faltan 12»."""
    dia = nueva(prs)
    cabecera(dia, "DESPACHO", "El sistema sabe qué había que evaluar, no solo qué se evaluó",
             "Sin un plan cargado, ningún tablero puede decir qué falta. Una ruta es ese plan.")
    tarjetas(dia, [
        ("Se arma en el panel",
         "El coordinador escribe las direcciones y se las asigna a un ingeniero. "
         "Solo aparecen las matrículas vigentes de su brigada."),
        ("Baja al teléfono y funciona sin señal",
         "Se descarga al sincronizar y se trabaja con el aparato desconectado, que es "
         "como se trabaja en terreno después de un sismo."),
        ("La visita se cierra sola",
         "Al enviar la evaluación, el servidor marca esa parada como hecha. Nadie "
         "transcribe una lista de cumplimiento al final del día."),
        ("Lo que no se pudo, también cuenta",
         "Nadie atendió, la dirección no existe, se negaron, no se pudo acceder. "
         "Cuatro botones grandes, sin escribir, sin señal."),
        ("Cobertura real, sin inflarla",
         "El avance se mide contra lo planeado. Lo levantado fuera de ruta se cuenta "
         "aparte: sumarlo daría 100 % a quien evaluó los predios equivocados."),
        ("La lista de direcciones se borra sola",
         "Es dato personal de gente que todavía no ha sido visitada. Vence en el "
         "teléfono por su propio reloj, sin depender de que haya señal."),
    ])
    pie_contacto(dia)
    return dia


def ventajas(prs):
    dia = nueva(prs)
    cabecera(dia, "VENTAJAS", "Por qué esta herramienta y no una planilla")
    tarjetas(dia, [
        ("Funciona sin señal",
         "No es un modo degradado: es el modo normal de operación. Se diseñó para el peor "
         "escenario de conectividad, no para el mejor."),
        ("Criterio uniforme",
         "Todas las brigadas aplican la misma regla. El consolidado es comparable entre "
         "municipios y entre equipos."),
        ("Trazabilidad completa",
         "Cada evaluación queda con matrícula, brigada, fecha, coordenada y fotos. Si el "
         "criterio se modificó, queda el motivo escrito."),
        ("Solo firma quien puede",
         "Sin matrícula la aplicación no permite guardar. Y ningún desalojo queda con una "
         "sola firma: cada uno pasa por una segunda revisión."),
        ("Protección de datos",
         "Lo que sale hacia autoridades va agregado por sector, con umbral mínimo. El predio "
         "nunca se expone (Ley 1581 de 2012)."),
        # La respuesta directa a «el papel se pierde», que es el problema que plantea la
        # lámina 2 y que hasta ahora ninguna ventaja contestaba. El respaldo diario se
        # mudó a la lámina de despliegue, que es donde una entidad pregunta por eso.
        ("El trabajo no se pierde",
         "Guardado duradero en el teléfono, acuse del servidor por cada evaluación, y "
         "nada se borra del equipo al enviarlo."),
    ])
    pie_contacto(dia)
    return dia


def despliegue(prs):
    dia = nueva(prs)
    cabecera(dia, "CÓMO SE DESPLIEGA", "Dos formas, y la decisión no es irreversible")
    opciones = [
        (AZUL, "Sobre nuestro servidor",
         ["Nosotros alojamos el receptor y la base de datos.",
          "La entidad recibe accesos y empieza a operar el mismo día.",
          "Sin infraestructura, sin personal de sistemas, sin licencias.",
          "Cada evaluación queda atribuida a la brigada que la envió.",
          "Respaldo diario cifrado a otro servidor, probado restaurando."],
         "Recomendado para empezar, o para una emergencia en curso."),
        (AZUL_OSC, "Sobre servidor propio",
         ["La entidad instala el sistema en su propia infraestructura.",
          "Los datos nunca salen de sus servidores.",
          "El código es abierto y la instalación está documentada.",
          "Requiere un servidor con HTTPS y alguien que lo administre.",
          "Nosotros acompañamos la puesta en marcha."],
         "Recomendado cuando la política de datos lo exige."),
    ]
    x0, w = Inches(0.9), Inches(5.6)
    for i, (color, titulo, puntos, nota) in enumerate(opciones):
        x = Emu(x0 + i * (w + Inches(0.3)))
        caja(dia, x, Inches(2.5), w, Inches(3.75), relleno=BLANCO, borde=LINEA)
        caja(dia, x, Inches(2.5), w, Inches(0.62), relleno=color, borde=None)
        texto(dia, Emu(x + Inches(0.35)), Inches(2.68), Emu(w - Inches(0.7)), Inches(0.35),
              titulo, tam=17, negrita=True, color=BLANCO)
        texto(dia, Emu(x + Inches(0.35)), Inches(3.32), Emu(w - Inches(0.7)), Inches(2.0),
              [f"·  {p}" for p in puntos], tam=12.5, color=TINTA2, interlineado=1.45)
        texto(dia, Emu(x + Inches(0.35)), Inches(5.62), Emu(w - Inches(0.7)), Inches(0.4),
              nota, tam=12, negrita=True, color=color)
    texto(dia, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.35),
          "Se puede empezar en nuestro servidor y migrar a infraestructura propia después: "
          "es el mismo software y los datos son exportables.",
          tam=12.5, color=TENUE)
    pie_contacto(dia)
    return dia


def integracion(prs):
    dia = nueva(prs)
    cabecera(dia, "INTEGRACIÓN", "Los reportes llegan a los sistemas que ya usan",
             "API de consulta en funcionamiento: solo lectura, credenciales propias "
             "y alcance por municipio.")
    tarjetas(dia, [
        ("Formulario V2F del IDIGER",
         "Cada evaluación aplanada a las casillas del formulario oficial, con sus "
         "códigos. Es el estándar que ya usan las brigadas."),
        ("Consolidado por sector",
         "Conteos por municipio y barrio con el umbral de anonimato aplicado. Sin dato "
         "personal: apto para un tablero público."),
        ("Capa geográfica",
         "Los mismos datos como GeoJSON, listos para un visor municipal o un geoportal, "
         "sin llave de terceros ni cuenta comercial."),
        ("Detalle para la entidad dueña",
         "Direcciones y coordenadas, solo para quien responde por esos datos y con "
         "finalidad declarada (Ley 1581 de 2012)."),
    ], y=Inches(2.85), alto=Inches(1.9), cols=4)
    caja(dia, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.3), relleno=PAPEL, borde=LINEA)
    texto(dia, Inches(1.2), Inches(5.45), Inches(10.9), Inches(0.9),
          "Las credenciales de lectura están separadas de las de brigada: una brigada "
          "escribe desde un teléfono, un sistema lee desde la alcaldía. Si fueran la misma, "
          "filtrar el token de un geoportal permitiría escribir evaluaciones falsas. Cada "
          "credencial se limita a sus municipios y se revoca en el acto.",
          tam=12.5, color=TINTA2, interlineado=1.3)
    pie_contacto(dia)
    return dia


def libre(prs):
    dia = nueva(prs)
    cabecera(dia, "LICENCIA", "Software libre: qué gana una entidad pública")
    tarjetas(dia, [
        ("Sin dependencia del proveedor",
         "El código es público y auditable. Si mañana esta empresa desaparece, la herramienta "
         "sigue funcionando y cualquiera puede mantenerla."),
        ("Sin costo por usuario",
         "No hay licencias por inspector ni por brigada. Escalar de diez a doscientos "
         "evaluadores no cambia el costo del software."),
        ("Adaptable a su normativa",
         "La regla de clasificación está aislada y documentada. Adaptarla a otro criterio o "
         "a otro país es una modificación acotada."),
        ("Las mejoras vuelven al común",
         "La licencia AGPL v3 obliga a que quien modifique y despliegue publique sus cambios. "
         "En emergencias, nadie debería reescribir esto desde cero."),
        ("Auditable por terceros",
         "Una universidad o un ente de control puede revisar exactamente cómo se calcula "
         "cada clasificación. No hay caja negra."),
        ("Documentación incluida",
         "Manual de campo, de coordinación y de operación del servidor, bajo la misma "
         "licencia que el programa."),
    ])
    pie_contacto(dia)
    return dia


def para_quien(prs):
    dia = nueva(prs)
    cabecera(dia, "QUÉ APORTA A CADA UNO", "Tres instituciones, tres necesidades distintas")
    bloques = [
        ("Alcaldías y gestión del riesgo",
         ["Cobertura verificable de la zona",
          "Consolidado por barrio para la crisis",
          "Trazabilidad de cada decisión",
          "Protección de datos por diseño",
          "Operación inmediata, sin compra"]),
        ("Universidades",
         ["Sus brigadas operan con criterio uniforme",
          "Estudiantes documentan sin firmar",
          "Registro del trabajo, con autoría",
          "Código abierto como caso de estudio",
          "Datos propios para investigación"]),
        ("Asociaciones profesionales",
         ["Control de qué matrículas están habilitadas",
          "Respaldo documental del ejercicio",
          "Criterio técnico común entre agremiados",
          "Evidencia ante revisión disciplinaria",
          "Instrumento de formación"]),
    ]
    x0, w, gap = Inches(0.9), Inches(3.65), Inches(0.28)
    for i, (titulo, puntos) in enumerate(bloques):
        x = Emu(x0 + i * (w + gap))
        caja(dia, x, Inches(2.3), w, Inches(4.1), relleno=BLANCO, borde=LINEA)
        caja(dia, x, Inches(2.3), w, Pt(5), relleno=AZUL, borde=None, radio=False)
        texto(dia, Emu(x + Inches(0.32)), Inches(2.62), Emu(w - Inches(0.64)), Inches(0.6),
              titulo, tam=15, negrita=True, color=AZUL_OSC, interlineado=1.15)
        texto(dia, Emu(x + Inches(0.32)), Inches(3.35), Emu(w - Inches(0.64)), Inches(2.9),
              [f"·  {p}" for p in puntos], tam=11.5, color=TINTA2, interlineado=1.45)
    pie_contacto(dia)
    return dia


def empezar(prs):
    dia = nueva(prs)
    cabecera(dia, "CÓMO EMPEZAR", "Lo que hace falta para tener una brigada operando")
    pasos = [
        ("Una conversación", "Entender su operación, su normativa interna y qué necesita el "
         "consolidado que ustedes ya producen."),
        ("Una prueba real", "Una brigada, un sector, una jornada. Sin compromiso y sin "
         "instalar nada en su infraestructura."),
        # Estas tarjetas son angostas: cuatro líneas y se sale. Enumerar los tres manuales
        # no cabía, así que va la idea —uno por rol— y el detalle queda para la landing.
        ("Ajuste y capacitación", "Adaptar lo que haga falta y formar a los inspectores, "
         "con un manual para cada rol."),
        ("Decidir dónde viven los datos", "Seguir en nuestro servidor o migrar a "
         "infraestructura propia. La decisión no es irreversible."),
    ]
    x0, w, gap = Inches(0.9), Inches(2.71), Inches(0.22)
    for i, (titulo, cuerpo) in enumerate(pasos):
        x = Emu(x0 + i * (w + gap))
        caja(dia, x, Inches(2.6), w, Inches(2.5), relleno=PAPEL, borde=LINEA)
        texto(dia, Emu(x + Inches(0.28)), Inches(2.85), Emu(w - Inches(0.56)), Inches(0.3),
              f"0{i+1}", tam=13, negrita=True, color=AZUL)
        texto(dia, Emu(x + Inches(0.28)), Inches(3.25), Emu(w - Inches(0.56)), Inches(0.6),
              titulo, tam=15, negrita=True, color=TINTA, interlineado=1.15)
        texto(dia, Emu(x + Inches(0.28)), Inches(4.0), Emu(w - Inches(0.56)), Inches(1.0),
              cuerpo, tam=12, color=TINTA2, interlineado=1.3)
    caja(dia, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.05), relleno=BLANCO, borde=AZUL,
         grosor=1.5)
    texto(dia, Inches(1.2), Inches(5.6), Inches(10.9), Inches(0.65),
          ["La aplicación ya está en línea y se puede probar hoy mismo desde cualquier teléfono.",
           "Escríbanos y coordinamos una demostración con su equipo."],
          tam=14, color=TINTA, interlineado=1.35)
    pie_contacto(dia)
    return dia


def contacto(prs):
    dia = nueva(prs, fondo=TINTA)
    caja(dia, Emu(0), Emu(0), Inches(0.34), ALTO, relleno=AZUL, borde=None, radio=False)
    texto(dia, Inches(1.3), Inches(1.5), Inches(10.5), Inches(0.4), "CONTACTO",
          tam=13, color=RGBColor(0x7D, 0xD3, 0xFC), negrita=True)
    texto(dia, Inches(1.3), Inches(2.05), Inches(10.5), Inches(1.0),
          "Hablemos de su operación", tam=40, color=BLANCO, negrita=True)
    barra_semaforo(dia, Inches(1.3), Inches(3.3), Inches(4.2))
    texto(dia, Inches(1.3), Inches(3.95), Inches(6.0), Inches(1.8), [
        (CONTACTO["nombre"], {"tam": 20, "negrita": True, "color": BLANCO}),
        (CONTACTO["empresa"], {"tam": 14, "color": RGBColor(0x94, 0xA3, 0xB8),
                               "espacio_antes": 4}),
    ], interlineado=1.3)
    texto(dia, Inches(1.3), Inches(5.15), Inches(7.0), Inches(1.2), [
        (CONTACTO["correo"], {"tam": 17, "color": RGBColor(0x7D, 0xD3, 0xFC), "negrita": True}),
        (CONTACTO["web"], {"tam": 17, "color": RGBColor(0x7D, 0xD3, 0xFC),
                           "espacio_antes": 8}),
    ], interlineado=1.3)
    caja(dia, Inches(8.3), Inches(3.9), Inches(4.0), Inches(2.4),
         relleno=RGBColor(0x1E, 0x29, 0x3B), borde=RGBColor(0x33, 0x41, 0x55))
    texto(dia, Inches(8.65), Inches(4.2), Inches(3.3), Inches(1.9), [
        ("Pruébela ahora", {"tam": 15, "negrita": True, "color": BLANCO}),
        ("brigadaestructural.co", {"tam": 14, "color": RGBColor(0x7D, 0xD3, 0xFC),
                                   "espacio_antes": 8}),
        ("Código fuente y documentación:", {"tam": 11, "color": RGBColor(0x94, 0xA3, 0xB8),
                                            "espacio_antes": 14}),
        ("github.com/abenito32/brigadaestructural", {"tam": 11,
                                                     "color": RGBColor(0x94, 0xA3, 0xB8)}),
    ], interlineado=1.25)
    texto(dia, Inches(1.3), Inches(6.85), Inches(10.8), Inches(0.3),
          "Software libre bajo GNU AGPL v3  ·  Desarrollada con Amor por "
          + CONTACTO["nombre"], tam=10, color=RGBColor(0x64, 0x74, 0x8B))
    return dia


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = ANCHO, ALTO
    for constructor in (portada, problema, que_es, como_funciona, semaforo, panel, despacho, ventajas,
                        despliegue, integracion, libre, para_quien, empezar, contacto):
        constructor(prs)
    prs.save(SALIDA)
    kb = SALIDA.stat().st_size / 1024
    print(f"{SALIDA}  ({len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas, {kb:.0f} KB)")


if __name__ == "__main__":
    main()
